import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
import xlrd
import subprocess
import os


def extract_text(file_path: str) -> str:
    """Extract text from a document. Supports PDF, DOC, DOCX, PPT, PPTX, XLS, XLSX, TXT."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".doc":
        return _extract_doc(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".ppt":
        return _extract_ppt(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext == ".xls":
        return _extract_xls(file_path)
    elif ext == ".xlsx":
        return _extract_xlsx(file_path)
    elif ext == ".txt":
        return _extract_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(path: str) -> str:
    doc = fitz.open(path)
    text_parts = []
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "\n".join(text_parts).strip()


def _extract_docx(path: str) -> str:
    doc = Document(path)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    return "\n".join(text_parts).strip()


def _extract_doc(path: str) -> str:
    result = subprocess.run(["antiword", path], capture_output=True, text=True)
    if result.returncode != 0:
        raise ValueError(f"antiword failed: {result.stderr.strip()}")
    return result.stdout.strip()


def _extract_ppt(path: str) -> str:
    """Extract text from legacy .ppt files by converting to .pptx with LibreOffice."""
    import tempfile
    with tempfile.TemporaryDirectory() as tmpdir:
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pptx", "--outdir", tmpdir, path],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode != 0:
            raise ValueError(f"LibreOffice conversion failed: {result.stderr.strip()}")
        converted = os.path.join(tmpdir, os.path.splitext(os.path.basename(path))[0] + ".pptx")
        if not os.path.exists(converted):
            raise ValueError("LibreOffice conversion produced no output file")
        return _extract_pptx(converted)


def _extract_pptx(path: str) -> str:
    prs = Presentation(path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
    return "\n".join(text_parts).strip()


def _extract_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read().strip()


def _extract_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    text_parts = []
    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None and str(cell).strip():
                    text_parts.append(str(cell))
    wb.close()
    return "\n".join(text_parts).strip()


def _extract_xls(path: str) -> str:
    wb = xlrd.open_workbook(path)
    text_parts = []
    for sheet in wb.sheets():
        for row_idx in range(sheet.nrows):
            for col_idx in range(sheet.ncols):
                cell = sheet.cell(row_idx, col_idx)
                if cell.value and str(cell.value).strip():
                    text_parts.append(str(cell.value))
    return "\n".join(text_parts).strip()


def is_scanned_pdf(text: str, threshold: int = 50) -> bool:
    """Heuristic: if extracted text is very short, likely scanned/image PDF."""
    return len(text.strip()) < threshold
